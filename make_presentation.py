"""
Academic presentation: Active Vibration Control of an ASML Measurement Plate
PPF (Phase I) + DVF (Phase II)
Style: academic-pptx-skill guidelines (action titles, white bg, navy/blue, one exhibit per results slide)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# --------------------------------------------------------------------------
# Colours
NAVY      = RGBColor(0x1F, 0x4E, 0x79)
BLUE      = RGBColor(0x2E, 0x75, 0xB6)
ORANGE    = RGBColor(0xC0, 0x55, 0x0A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
DARK      = RGBColor(0x26, 0x26, 0x26)
GRAY      = RGBColor(0x59, 0x59, 0x59)
LGRAY     = RGBColor(0xD6, 0xD6, 0xD6)

W   = Inches(13.333)
H   = Inches(7.5)
MRG = Inches(0.55)   # margin

# --------------------------------------------------------------------------
# Helpers
def set_font(run, bold=False, size=None, colour=None, italic=False):
    run.font.bold   = bold
    run.font.italic = italic
    if size:    run.font.size  = Pt(size)
    if colour:  run.font.color.rgb = colour

def add_textbox(slide, text, x, y, w, h, font_size=20, bold=False,
                colour=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txb  = slide.shapes.add_textbox(x, y, w, h)
    tf   = txb.text_frame
    tf.word_wrap = wrap
    p    = tf.paragraphs[0]
    p.alignment = align
    run  = p.add_run()
    run.text = text
    set_font(run, bold=bold, size=font_size, colour=colour)
    return txb

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_w=Pt(1)):
    from pptx.util import Pt as Pt2
    shp = slide.shapes.add_shape(
        1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def slide_header(slide, action_title, section=None):
    """Add navy action title bar at top of slide."""
    # title bar background
    add_rect(slide, 0, 0, W, Inches(1.05), NAVY)
    # action title text
    add_textbox(slide, action_title,
                MRG, Inches(0.12), W - 2*MRG, Inches(0.85),
                font_size=22, bold=True, colour=WHITE, align=PP_ALIGN.LEFT)
    if section:
        add_textbox(slide, section,
                    W - Inches(3.2), Inches(0.78), Inches(3.0), Inches(0.25),
                    font_size=12, colour=RGBColor(0xBF, 0xD4, 0xED), align=PP_ALIGN.RIGHT)

def add_bullet_block(slide, bullets, x, y, w, h,
                     size=19, title=None, title_colour=BLUE):
    """Add a list of bullet strings as a text frame."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = title
        set_font(r, bold=True, size=size+1, colour=title_colour)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        first = False
        p.level = b.get('level', 0)
        r = p.add_run()
        r.text = b['text']
        set_font(r, bold=b.get('bold', False),
                 size=b.get('size', size),
                 colour=b.get('colour', DARK))

def add_callout_box(slide, text, x, y, w, h,
                    bg=RGBColor(0xE2, 0xEF, 0xF9), border=BLUE):
    """Highlight box for key finding."""
    add_rect(slide, x, y, w, h, bg, border, Pt(1.5))
    txb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.1),
                                   w - Inches(0.24), h - Inches(0.18))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    set_font(r, size=17, colour=NAVY)

def footer(slide, text="Isaak Bouwmeester | ASML Internship | June 2026",
           slide_num=None, total=None):
    add_rect(slide, 0, H - Inches(0.30), W, Inches(0.30), LGRAY)
    add_textbox(slide, text,
                MRG, H - Inches(0.28), Inches(8), Inches(0.26),
                font_size=11, colour=GRAY)
    if slide_num is not None:
        num = f"{slide_num}/{total}" if total else str(slide_num)
        add_textbox(slide, num,
                    W - Inches(1.2), H - Inches(0.28), Inches(1.0), Inches(0.26),
                    font_size=11, colour=GRAY, align=PP_ALIGN.RIGHT)

def divider_slide(prs, label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_textbox(slide, label,
                Inches(1), Inches(2.5), W - Inches(2), Inches(2.5),
                font_size=40, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    return slide

# --------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]
total_slides = 18

sn = 0  # slide number counter

# ==============================
# SLIDE 1: Title
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = NAVY

add_textbox(sl,
    "Active Vibration Control of an ASML Measurement Plate",
    Inches(0.8), Inches(1.6), W - Inches(1.6), Inches(1.8),
    font_size=36, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl,
    "Phase I: Positive Position Feedback with Dual Piezo Patches\n"
    "Phase II: Direct Velocity Feedback with MEMS Accelerometer",
    Inches(1.2), Inches(3.5), W - Inches(2.4), Inches(1.4),
    font_size=22, colour=RGBColor(0xBF, 0xD4, 0xED), align=PP_ALIGN.CENTER)

add_textbox(sl,
    "Isaak Bouwmeester | ASML Internship | June 2026",
    Inches(1.2), Inches(5.4), W - Inches(2.4), Inches(0.5),
    font_size=18, colour=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.CENTER)

# ==============================
# SLIDE 2: Problem
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "Floor vibrations at the plate's first bending mode degrade wafer positioning accuracy",
    "Motivation")
footer(sl, slide_num=sn, total=total_slides)

add_bullet_block(sl, [
    {'text': 'ASML measurement plate: ~1000x1000 mm aluminium, ~100 mm thick'},
    {'text': 'First bending mode at ~160 Hz (COMSOL); resonance amplifies floor vibrations'},
    {'text': 'Positioning errors at wafer reduce lithography resolution'},
    {'text': ''},
    {'text': 'Goal: suppress the dominant bending mode using active vibration control (AVC)'},
    {'text': 'Approach: piezoelectric patch actuator in closed-loop feedback'},
],
x=MRG, y=Inches(1.25), w=Inches(6.5), h=Inches(5.5))

add_callout_box(sl,
    "Two phases:\n"
    "  Phase I  - PPF with dual DuraAct patches (Patch 2 actuator + Patch 3 sensor)\n"
    "  Phase II - DVF with Patch 2 actuator + MEMS accelerometer (IMU)",
    Inches(7.2), Inches(1.4), Inches(5.8), Inches(1.5))

# Physical setup schematic (text-based)
add_rect(sl, Inches(7.2), Inches(3.2), Inches(5.8), Inches(3.1), OFF_WHITE, LGRAY)
add_textbox(sl, "Physical Setup",
            Inches(7.4), Inches(3.25), Inches(5.4), Inches(0.35),
            font_size=14, bold=True, colour=NAVY)
add_bullet_block(sl, [
    {'text': 'Plate:    Al, 1x1 m, 100 mm thick, STL import in COMSOL'},
    {'text': 'Actuator: PI Ceramic P-876.A12 DuraAct (d31 mode)'},
    {'text': '          3-layer: polymer / PIC255 ceramic / polymer'},
    {'text': 'Mounting: Domain Spring Foundation, f_mount ~ 1.7 Hz'},
    {'text': 'BC:       elastic isolation (NOT fixed constraint)'},
],
x=Inches(7.4), y=Inches(3.65), w=Inches(5.5), h=Inches(2.5), size=15)

# ==============================
# SLIDE 3: Eigenfrequencies
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "The plate has four flexible bending modes between 87 and 226 Hz; mode 3 at 161 Hz is the DVF target",
    "COMSOL Model")
footer(sl, slide_num=sn, total=total_slides)

# Table of eigenfrequencies
headers = ["Mode group", "Freq. (Hz)", "Description", "DVF phase at peak"]
rows = [
    ("Spring rigid-body (x6)", "1.66-1.68", "Translation + rotation on foam mount", "N/A"),
    ("Flexible mode 1", "86.51", "First bending (actuator placed here)", "+51 to +140 deg - distorted"),
    ("Flexible mode 2", "127.78", "Second bending", "~-90 deg"),
    ("Flexible mode 3", "160.90", "Third bending - DVF target", "-90 deg - clean"),
    ("Flexible mode 4", "225.56", "Fourth bending", "~-90 deg"),
]
col_w = [Inches(2.8), Inches(1.4), Inches(4.5), Inches(3.8)]
col_x = [MRG, MRG + Inches(2.9), MRG + Inches(4.4), MRG + Inches(9.0)]
row_h = Inches(0.54)
y0 = Inches(1.25)

# header row
add_rect(sl, MRG, y0, sum(col_w) + Inches(0.1), row_h, NAVY)
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_textbox(sl, hdr, cx + Inches(0.06), y0 + Inches(0.10),
                cw - Inches(0.1), row_h - Inches(0.12),
                font_size=15, bold=True, colour=WHITE)

for i, row in enumerate(rows):
    ry = y0 + (i+1) * row_h
    bg = RGBColor(0xE2, 0xEF, 0xF9) if i % 2 == 0 else WHITE
    # highlight the DVF target row
    if i == 3:
        bg = RGBColor(0xFF, 0xF2, 0xCC)
    add_rect(sl, MRG, ry, sum(col_w) + Inches(0.1), row_h, bg, LGRAY, Pt(0.5))
    for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        col_c = DARK
        if i == 3: col_c = RGBColor(0x7F, 0x3D, 0x00)
        add_textbox(sl, cell, cx + Inches(0.06), ry + Inches(0.10),
                    cw - Inches(0.1), row_h - Inches(0.12),
                    font_size=14, colour=col_c)

add_callout_box(sl,
    "Key: Mode 3 at 161 Hz shows a clean -90 deg phase at the |H22| peak. Mode 1 at 87 Hz "
    "has feedthrough-distorted phase (+51 to +140 deg) -- simple DVF cannot reliably suppress it.",
    MRG, Inches(4.95), W - 2*MRG, Inches(1.25))

# ==============================
# SLIDE 4: Actuator
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "The DuraAct P-876.A12 patch uses d31 bending; three COMSOL material overrides are required",
    "Actuator")
footer(sl, slide_num=sn, total=total_slides)

add_bullet_block(sl, [
    {'text': 'DuraAct layer stack (d31 bending mode, no rotation needed):'},
    {'text': 'Bottom polymer:  61x35 mm, 0.15 mm thick  (E = 5 GPa)', 'level': 1},
    {'text': 'PIC255 ceramic:  50x30 mm, 0.20 mm thick  (active layer)', 'level': 1},
    {'text': 'Top polymer:     61x35 mm, 0.15 mm thick', 'level': 1},
    {'text': ''},
    {'text': 'Electric field applied through ceramic thickness (z-axis = poling direction)'},
    {'text': 'Va = 1 V on top PIC255 face, Ground on bottom PIC255 face'},
],
x=MRG, y=Inches(1.25), w=Inches(6.3), h=Inches(3.8), size=19)

add_callout_box(sl,
    "CRITICAL: COMSOL PIC255 defaults are wrong. Always override:\n"
    "  d31 = d32 = -260 pm/V  (COMSOL default: -18.5e-11 C/N)\n"
    "  epsR_T33  = 1355        (COMSOL default: 1832)\n"
    "  Ceramic thickness = 0.20 mm",
    MRG, Inches(5.2), Inches(6.3), Inches(1.8))

# Layer diagram (text box schematic)
add_rect(sl, Inches(7.5), Inches(1.4), Inches(5.5), Inches(4.8), OFF_WHITE, LGRAY)
add_textbox(sl, "Patch cross-section (not to scale)",
            Inches(7.7), Inches(1.45), Inches(5.1), Inches(0.35),
            font_size=14, bold=True, colour=NAVY)
layers = [
    ("Top polymer   0.15 mm",   RGBColor(0xC0, 0xD4, 0xE9)),
    ("PIC255 ceramic  0.20 mm", RGBColor(0xF4, 0xB1, 0x83)),
    ("Bottom polymer 0.15 mm",  RGBColor(0xC0, 0xD4, 0xE9)),
    ("Plate (Al)  100 mm",      RGBColor(0xD9, 0xD9, 0xD9)),
]
lay_h = [Inches(0.55), Inches(0.65), Inches(0.55), Inches(1.5)]
ly = Inches(1.9)
for (lbl, rgb), lh in zip(layers, lay_h):
    add_rect(sl, Inches(7.8), ly, Inches(5.0), lh, rgb, DARK, Pt(0.5))
    add_textbox(sl, lbl, Inches(7.9), ly + Inches(0.1),
                Inches(4.8), lh - Inches(0.12), font_size=15, colour=DARK)
    ly += lh

# ==============================
# PART I DIVIDER
# ==============================
sn += 1
divider_slide(prs, "Part I\nPositive Position Feedback\nwith Dual Piezo Patches")

# ==============================
# SLIDE 5: PPF setup
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "Phase I bonds a second DuraAct patch as an open-circuit voltage sensor for PPF feedback",
    "Phase I Setup")
footer(sl, slide_num=sn, total=total_slides)

add_bullet_block(sl, [
    {'text': 'Patch 3 (sensor): identical DuraAct layup, offset 65 mm in x from Patch 2'},
    {'text': 'Open-circuit condition: Floating Potential on both ceramic faces (top + bottom)'},
    {'text': '  Do NOT add Ground to either face -- that fixes V = 0 and corrupts Vs', 'level': 1},
    {'text': 'Sensor output: Vs = open-circuit piezo voltage  [V/V]'},
    {'text': ''},
    {'text': 'COMSOL model: 7 domains (plate + 3 Patch-2 layers + 3 Patch-3 layers)'},
    {'text': 'Boundary condition: Fixed Constraint on 4 bottom perimeter edge lines'},
    {'text': '  Geometric entity level must be set to Edge, not Boundary', 'level': 1},
],
x=MRG, y=Inches(1.25), w=Inches(6.3), h=Inches(5.5), size=19)

# Domain diagram
add_rect(sl, Inches(7.3), Inches(1.4), Inches(5.7), Inches(5.5), OFF_WHITE, LGRAY)
add_textbox(sl, "COMSOL domain layout",
            Inches(7.5), Inches(1.45), Inches(5.3), Inches(0.35),
            font_size=14, bold=True, colour=NAVY)

domains = [
    ("Plate (aluminium)", RGBColor(0xD9, 0xD9, 0xD9), 1.9, 2.3),
    ("Patch 2 - top polymer", RGBColor(0xC0, 0xD4, 0xE9), 4.3, 0.45),
    ("Patch 2 - PIC255 ceramic", RGBColor(0xF4, 0xB1, 0x83), 4.8, 0.5),
    ("Patch 2 - bottom polymer", RGBColor(0xC0, 0xD4, 0xE9), 5.35, 0.45),
    ("Patch 3 - top polymer", RGBColor(0xC0, 0xD4, 0xE9), 5.85, 0.45),
    ("Patch 3 - PIC255 ceramic", RGBColor(0xF4, 0xB1, 0x83), 6.35, 0.5),
    ("Patch 3 - bottom polymer", RGBColor(0xC0, 0xD4, 0xE9), 6.9, 0.45),
]
for name, rgb, iny, inh in domains:
    add_rect(sl, Inches(7.8), Inches(iny), Inches(4.9), Inches(inh),
             rgb, DARK, Pt(0.5))
    add_textbox(sl, name, Inches(7.85), Inches(iny + 0.07),
                Inches(4.8), Inches(inh), font_size=13, colour=DARK)

# ==============================
# SLIDE 6: PPF theory + block diagram
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "PPF is unconditionally stable because the collocated piezo phase is always within (-90 deg, +90 deg)",
    "PPF Theory")
footer(sl, slide_num=sn, total=total_slides)

add_bullet_block(sl, [
    {'text': 'Plant:   H22(s) = H0 * wn^2 / (s^2 + 2*zeta*wn*s + wn^2)'},
    {'text': 'Phase of H22 at resonance: -90 deg  (collocated piezo pair)'},
    {'text': ''},
    {'text': 'PPF controller (Baken 2025):'},
    {'text': '  K_PPF(s) = (1/H0) * g * wf^2 / (s^2 + 2*zf*wf*s + wf^2)', 'level': 1},
    {'text': '  wf = wn,  g = 0.3,  zf = 0.3  (Baken values)', 'level': 1},
    {'text': ''},
    {'text': 'Phase accumulation at resonance:'},
    {'text': '  H22: displacement to Vs = -90 deg', 'level': 1},
    {'text': '  K_PPF filter: Vs to Va  = -90 deg', 'level': 1},
    {'text': '  Total: -180 deg from displacement = opposing velocity = DAMPING', 'level': 1, 'bold': True},
],
x=MRG, y=Inches(1.25), w=Inches(6.5), h=Inches(5.6), size=18)

# Block diagram (drawn with shapes)
bx, by = Inches(7.6), Inches(1.6)
bw = Inches(5.4)

# PPF block diagram
add_textbox(sl, "PPF block diagram:",
            bx, by, bw, Inches(0.35), font_size=15, bold=True, colour=NAVY)

# summing junction circle
cir = sl.shapes.add_shape(9, bx + Inches(0.4), by + Inches(0.55), Inches(0.55), Inches(0.55))
cir.fill.solid(); cir.fill.fore_color.rgb = WHITE
cir.line.color.rgb = DARK; cir.line.width = Pt(2)
add_textbox(sl, "+\n+", bx + Inches(0.46), by + Inches(0.60),
            Inches(0.45), Inches(0.5), font_size=13, colour=DARK, align=PP_ALIGN.CENTER)

# d_dis arrow label
add_textbox(sl, "d_dis ->", bx, by + Inches(0.67), Inches(0.42), Inches(0.3),
            font_size=13, colour=DARK)

# Va label + arrow to H22
add_textbox(sl, "Va", bx + Inches(1.1), by + Inches(0.67), Inches(0.35), Inches(0.3),
            font_size=13, colour=DARK)

# H22 box
h22b = sl.shapes.add_shape(1, bx + Inches(1.5), by + Inches(0.5), Inches(1.5), Inches(0.65))
h22b.fill.solid(); h22b.fill.fore_color.rgb = BLUE
h22b.line.color.rgb = NAVY; h22b.line.width = Pt(1.5)
add_textbox(sl, "H22", bx + Inches(1.6), by + Inches(0.57),
            Inches(1.3), Inches(0.5), font_size=18, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

# Vs label
add_textbox(sl, "-> Vs ->", bx + Inches(3.1), by + Inches(0.67), Inches(1.0), Inches(0.3),
            font_size=13, colour=DARK)

# K_PPF box
kppf = sl.shapes.add_shape(1, bx + Inches(1.5), by + Inches(1.6), Inches(1.5), Inches(0.65))
kppf.fill.solid(); kppf.fill.fore_color.rgb = ORANGE
kppf.line.color.rgb = RGBColor(0x80, 0x30, 0x00); kppf.line.width = Pt(1.5)
add_textbox(sl, "K_PPF", bx + Inches(1.6), by + Inches(1.67),
            Inches(1.3), Inches(0.5), font_size=18, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, "(feedback: positive +)",
            bx, by + Inches(2.55), bw, Inches(0.3), font_size=13,
            colour=GRAY, align=PP_ALIGN.CENTER)

# arrows (lines described in text since pptx connectors are complex)
add_textbox(sl,
    "[d_dis] -> (+) -> H22 -> [Vs] ->\n"
    "             ^                |\n"
    "             |   K_PPF  <----|",
    bx, by + Inches(2.95), bw, Inches(0.95), font_size=14,
    colour=NAVY)

add_callout_box(sl,
    "Stability: phase of H22 bounded in (-180 deg, 0 deg).\n"
    "PPF filter adds (-180 deg, 0 deg). Sum never encircles -1. -> Unconditionally stable.",
    bx, by + Inches(4.0), bw, Inches(1.1))

# ==============================
# SLIDE 7: PPF Closed-Loop
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "PPF predicts 22.6 dB suppression at f1 = 214 Hz with g = 0.3 -- matching Baken 2025",
    "PPF Results")
footer(sl, slide_num=sn, total=total_slides)

add_bullet_block(sl, [
    {'text': 'Closed-loop tip displacement (MIMO formula):'},
    {'text': 'T11 = H11 - H12 * [K_PPF / (1 - H22*K_PPF)] * H21', 'level': 1, 'bold': True},
    {'text': ''},
    {'text': 'MATLAB implementation:'},
    {'text': 'T = feedback(HOLc, K_ppf, 2, 2, +1)  (positive feedback, channel 2)', 'level': 1},
    {'text': ''},
    {'text': 'Tuning: g = 0.3,  zeta_f = 0.3  (Baken 2025 values)'},
    {'text': ''},
    {'text': 'Predicted suppression at f1 = 214 Hz:'},
    {'text': 'Delta_H22  = -22.6 dB', 'level': 1, 'bold': True},
    {'text': 'Baken 2025 result: -22.8 dB  (sandwich beam, same g/zeta_f)', 'level': 1},
    {'text': ''},
    {'text': 'Discretisation: Tustin method, fs = 10 kHz'},
    {'text': 'Kd = c2d(K_ppf, 1e-4, tustin)', 'level': 1},
],
x=MRG, y=Inches(1.25), w=Inches(6.5), h=Inches(5.6), size=18)

# Result comparison table
add_rect(sl, Inches(7.5), Inches(1.4), Inches(5.5), Inches(5.5), OFF_WHITE, LGRAY)
add_textbox(sl, "Comparison: this project vs Baken 2025",
            Inches(7.7), Inches(1.45), Inches(5.1), Inches(0.35),
            font_size=14, bold=True, colour=NAVY)
comp_rows = [
    ("Parameter", "This project", "Baken 2025"),
    ("Structure", "Al plate 1x1 m", "Sandwich beam 0.4 m"),
    ("Eigenfreq f1", "214 Hz", "29.5 Hz"),
    ("PPF gain g", "0.3", "0.3"),
    ("PPF damp. zf", "0.3", "0.3"),
    ("Suppression", "-22.6 dB", "-22.8 dB"),
    ("H22 at DC (H0)", "~0.22 V/V (est.)", "read from FRD"),
]
ry2 = Inches(1.9)
for i, row in enumerate(comp_rows):
    bg2 = NAVY if i == 0 else (RGBColor(0xE2, 0xEF, 0xF9) if i % 2 else WHITE)
    fc2 = WHITE if i == 0 else DARK
    add_rect(sl, Inches(7.5), ry2, Inches(5.5), Inches(0.5), bg2, LGRAY, Pt(0.5))
    col_xs = [Inches(7.6), Inches(9.5), Inches(11.3)]
    col_ws = [Inches(1.85), Inches(1.75), Inches(1.65)]
    for cell, cx2, cw2 in zip(row, col_xs, col_ws):
        add_textbox(sl, cell, cx2, ry2 + Inches(0.07),
                    cw2, Inches(0.38), font_size=13, colour=fc2, bold=(i==0))
    ry2 += Inches(0.5)

# ==============================
# PART II DIVIDER
# ==============================
sn += 1
divider_slide(prs, "Part II\nDirect Velocity Feedback\nwith MEMS Accelerometer")

# ==============================
# SLIDE 8: Phase convention
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "COMSOL's solid.accZ uses +omega^2 (not -omega^2), shifting the phase 180 deg vs a real accelerometer",
    "Phase Convention")
footer(sl, slide_num=sn, total=total_slides)

add_bullet_block(sl, [
    {'text': 'Physical acceleration: a = (j*omega)^2 * x = -omega^2 * x'},
    {'text': '  Phase = +180 deg relative to displacement', 'level': 1},
    {'text': ''},
    {'text': 'COMSOL (solid.accZ): derived from inertia term (+omega^2 * M * U)'},
    {'text': '  solid.accZ = +omega^2 * w  = -a_physical', 'level': 1},
    {'text': ''},
    {'text': 'Consequence: 180 deg phase offset vs Falangas 1994 and all papers using real accelerometers'},
],
x=MRG, y=Inches(1.25), w=Inches(6.5), h=Inches(3.5), size=19)

# Phase comparison table
headers2 = ["Freq. region", "solid.accZ phase", "Physical accel. phase"]
rows2 = [
    ("Below resonance (DC)", "0 deg", "+180 deg"),
    ("At resonance w1", "-90 deg", "-270 deg (+90 deg)"),
    ("Above resonance", "-180 deg", "-360 deg (0 deg)"),
]
ry3 = Inches(4.9)
add_rect(sl, MRG, ry3, W - 2*MRG, Inches(0.45), NAVY)
for j, h in enumerate(headers2):
    xs3 = [MRG + Inches(0.1), MRG + Inches(4.2), MRG + Inches(7.5)]
    add_textbox(sl, h, xs3[j], ry3 + Inches(0.07), Inches(3.9), Inches(0.33),
                font_size=15, bold=True, colour=WHITE)
ry3 += Inches(0.45)
for i, row in enumerate(rows2):
    bg3 = RGBColor(0xE2, 0xEF, 0xF9) if i % 2 == 0 else WHITE
    add_rect(sl, MRG, ry3, W - 2*MRG, Inches(0.45), bg3, LGRAY, Pt(0.5))
    for j, cell in enumerate(row):
        xs3 = [MRG + Inches(0.1), MRG + Inches(4.2), MRG + Inches(7.5)]
        add_textbox(sl, cell, xs3[j], ry3 + Inches(0.07), Inches(3.9), Inches(0.33),
                    font_size=15, colour=DARK)
    ry3 += Inches(0.45)

add_callout_box(sl,
    "Confirmed from H22_accel.txt data: at 1 Hz, H22 is positive real (0 deg) -- COMSOL convention. "
    "Resonance at ~157 Hz, phase crosses -90 deg there.",
    MRG, ry3 + Inches(0.05), W - 2*MRG, Inches(0.65))

# ==============================
# SLIDE 9: Controller Selection
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "PPF and IRC both fail with an acceleration sensor; DVF is conditionally stable and applicable",
    "Controller Selection")
footer(sl, slide_num=sn, total=total_slides)

ctrl_headers = ["Strategy", "Stability condition", "H22 phase at resonance", "Verdict"]
ctrl_rows = [
    ("PPF", "Angle H22 in (-90, +90) deg", "-90 deg -- on boundary", "REJECTED: unsafe"),
    ("IRC", "Im[H22] <= 0 for all w\n(Negative Imaginary)", "NI holds only for disp./\nstrain sensors, not accel.", "REJECTED: proof invalid"),
    ("DVF", "Nyquist: g*H22_vel does\nnot encircle -1", "H22_vel real and positive\nat resonance", "SELECTED: conditionally stable"),
]
col_ws2 = [Inches(1.8), Inches(3.2), Inches(3.2), Inches(4.3)]
col_xs2 = [MRG, MRG + Inches(1.9), MRG + Inches(5.2), MRG + Inches(8.5)]
ry4 = Inches(1.3)
add_rect(sl, MRG, ry4, W - 2*MRG, Inches(0.45), NAVY)
for j, (h, cx, cw) in enumerate(zip(ctrl_headers, col_xs2, col_ws2)):
    add_textbox(sl, h, cx + Inches(0.05), ry4 + Inches(0.07),
                cw - Inches(0.05), Inches(0.33), font_size=15, bold=True, colour=WHITE)
ry4 += Inches(0.45)
row_bgs = [RGBColor(0xFF, 0xE6, 0xE6), RGBColor(0xFF, 0xE6, 0xE6),
           RGBColor(0xE2, 0xF0, 0xD9)]
for i, (row, bg) in enumerate(zip(ctrl_rows, row_bgs)):
    rh = Inches(1.1)
    add_rect(sl, MRG, ry4, W - 2*MRG, rh, bg, LGRAY, Pt(0.5))
    for j, (cell, cx, cw) in enumerate(zip(row, col_xs2, col_ws2)):
        bold_v = (j == 3)
        add_textbox(sl, cell, cx + Inches(0.05), ry4 + Inches(0.08),
                    cw - Inches(0.08), rh - Inches(0.1),
                    font_size=15, colour=DARK, bold=bold_v)
    ry4 += rh

add_callout_box(sl,
    "DVF conditional stability (Gatti 2007): patch actuator applies a distributed bending moment, "
    "not a point force -- Balas (1979) unconditional stability does NOT apply. "
    "Maximum stable gain g_max must be found from the Nyquist plot.",
    MRG, ry4 + Inches(0.05), W - 2*MRG, Inches(0.9))

# ==============================
# SLIDE 10: DVF block diagram + sign
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "DVF uses +g sign (not -g) with solid.accZ data -- the literature sign assumes physical acceleration",
    "DVF Sign Convention")
footer(sl, slide_num=sn, total=total_slides)

add_bullet_block(sl, [
    {'text': 'DVF controller (Falangas 1994 formulation):'},
    {'text': '  K_DVF(s) = +g * S / (s + S),   S = 5 rad/s  (~0.8 Hz cutoff)', 'level': 1, 'bold': True},
    {'text': '  Feedback: NEGATIVE (summing junction carries minus sign)', 'level': 1},
    {'text': ''},
    {'text': 'Why +g (not -g from paper):'},
    {'text': 'Falangas writes K = -g*S/(s+S) assuming a = -omega^2 * x  (physical)', 'level': 1},
    {'text': 'solid.accZ uses a = +omega^2 * x  (COMSOL convention)', 'level': 1},
    {'text': 'Feeding solid.accZ into -g flips damping -> anti-damping', 'level': 1},
    {'text': '-> use +g to restore genuine velocity damping', 'level': 1, 'bold': True},
    {'text': ''},
    {'text': 'Symptom of wrong sign:'},
    {'text': 'Suppression dB is POSITIVE and grows with gain (amplification, not suppression)', 'level': 1},
    {'text': 'Nyquist gain margin lands exactly on f1  (should be a different mode)', 'level': 1},
],
x=MRG, y=Inches(1.25), w=Inches(6.8), h=Inches(5.6), size=18)

# DVF block diagram (right column)
bx2, by2 = Inches(7.8), Inches(1.5)
add_textbox(sl, "DVF block diagram:",
            bx2, by2, Inches(5.0), Inches(0.35), font_size=15, bold=True, colour=NAVY)

# summing junction
cir2 = sl.shapes.add_shape(9, bx2 + Inches(0.4), by2 + Inches(0.55), Inches(0.55), Inches(0.55))
cir2.fill.solid(); cir2.fill.fore_color.rgb = WHITE
cir2.line.color.rgb = DARK; cir2.line.width = Pt(2)
add_textbox(sl, "+\n-", bx2 + Inches(0.46), by2 + Inches(0.60),
            Inches(0.45), Inches(0.5), font_size=13, colour=DARK, align=PP_ALIGN.CENTER)
add_textbox(sl, "d_dis ->", bx2, by2 + Inches(0.67), Inches(0.42), Inches(0.3),
            font_size=13, colour=DARK)
add_textbox(sl, "Va", bx2 + Inches(1.1), by2 + Inches(0.67), Inches(0.35), Inches(0.3),
            font_size=13, colour=DARK)

h22b2 = sl.shapes.add_shape(1, bx2 + Inches(1.5), by2 + Inches(0.5), Inches(1.5), Inches(0.65))
h22b2.fill.solid(); h22b2.fill.fore_color.rgb = BLUE
h22b2.line.color.rgb = NAVY; h22b2.line.width = Pt(1.5)
add_textbox(sl, "H22", bx2 + Inches(1.6), by2 + Inches(0.57),
            Inches(1.3), Inches(0.5), font_size=18, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, "-> a_imu ->", bx2 + Inches(3.1), by2 + Inches(0.67),
            Inches(1.2), Inches(0.3), font_size=13, colour=DARK)

kdvf = sl.shapes.add_shape(1, bx2 + Inches(1.5), by2 + Inches(1.6), Inches(1.5), Inches(0.65))
kdvf.fill.solid(); kdvf.fill.fore_color.rgb = ORANGE
kdvf.line.color.rgb = RGBColor(0x80, 0x30, 0x00); kdvf.line.width = Pt(1.5)
add_textbox(sl, "+g*S/(s+S)", bx2 + Inches(1.52), by2 + Inches(1.67),
            Inches(1.46), Inches(0.5), font_size=14, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl,
    "[d_dis] -> (+/-) -> H22 -> [a_imu] ->\n"
    "               ^                   |\n"
    "               |   +g*S/(s+S)  <---|",
    bx2, by2 + Inches(2.55), Inches(5.0), Inches(0.95), font_size=14, colour=NAVY)

add_callout_box(sl,
    "Numerics: at resonance H22*K = (-jM)(j*gL) = +MgL > 0\n"
    "Denominator = 1 + MgL > 1  ->  T21 = H21/(1+MgL) < H21  -> genuine suppression",
    bx2, by2 + Inches(3.7), Inches(5.0), Inches(1.1))

# ==============================
# SLIDE 11: Boundary conditions
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "Only Domain Spring Foundation is consistent between Study 1 (Va) and Study 2 (base) for a correct MIMO matrix",
    "Boundary Conditions")
footer(sl, slide_num=sn, total=total_slides)

bc_headers = ["BC approach", "Study 1 (H22, H12)", "Study 2 (H11, H21)", "Verdict"]
bc_rows = [
    ("Fixed Constraint\n(4 bottom edges)",
     "Works: f1 ~214 Hz",
     "Over-constrains plate:\nf1 jumps to ~6500 Hz",
     "Phase I only"),
    ("Boundary Spring Foundation",
     "Works",
     "FAILS: contradicts\nBase Excitation node\n-> flat H11",
     "REJECTED"),
    ("Rigid Motion Suppression (RMS)",
     "Works",
     "Different eigenmodes\nthan Study 1 ->\ninconsistent MIMO",
     "REJECTED"),
    ("Domain Spring Foundation\n+ Base Excitation node",
     "Works: 6 spring modes\nat ~1.7 Hz",
     "Correct: H21 ~ 1 at DC,\nresonances at correct freq.",
     "SELECTED: Phase II"),
]
col_ws3 = [Inches(2.8), Inches(2.5), Inches(3.2), Inches(3.0)]
col_xs3 = [MRG, MRG + Inches(2.9), MRG + Inches(5.5), MRG + Inches(8.8)]
ry5 = Inches(1.3)
add_rect(sl, MRG, ry5, W - 2*MRG, Inches(0.45), NAVY)
for j, (h, cx, cw) in enumerate(zip(bc_headers, col_xs3, col_ws3)):
    add_textbox(sl, h, cx + Inches(0.05), ry5 + Inches(0.07),
                cw - Inches(0.05), Inches(0.33), font_size=14, bold=True, colour=WHITE)
ry5 += Inches(0.45)
bc_bgs = [WHITE, RGBColor(0xFF, 0xE6, 0xE6), RGBColor(0xFF, 0xE6, 0xE6), RGBColor(0xE2, 0xF0, 0xD9)]
for i, (row, bg) in enumerate(zip(bc_rows, bc_bgs)):
    rh = Inches(1.0)
    add_rect(sl, MRG, ry5, W - 2*MRG, rh, bg, LGRAY, Pt(0.5))
    for j, (cell, cx, cw) in enumerate(zip(row, col_xs3, col_ws3)):
        add_textbox(sl, cell, cx + Inches(0.05), ry5 + Inches(0.07),
                    cw - Inches(0.08), rh - Inches(0.1), font_size=13, colour=DARK)
    ry5 += rh

# ==============================
# SLIDE 12: DVF results
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "DVF predicts 10-15 dB suppression at 161 Hz, consistent with Falangas 1994's 10.5 dB benchmark",
    "DVF Results")
footer(sl, slide_num=sn, total=total_slides)

add_bullet_block(sl, [
    {'text': 'Closed-loop transmissibility (primary metric: IMU corner):'},
    {'text': 'T21 = H21 / (1 + H22 * K_DVF)', 'level': 1, 'bold': True},
    {'text': ''},
    {'text': 'Controller: K_DVF = +g * 5/(s+5),  g = 0.30 * g_max'},
    {'text': ''},
    {'text': 'g_max from Nyquist: g_max = 1 / max(-Re[H22 * S/(jw+S)])'},
    {'text': 'Operating gain: g = 30% of g_max  (10 dB stability margin)'},
    {'text': ''},
    {'text': 'Predicted suppression: Delta_T21 ~ -10 to -15 dB at f1 = 161 Hz'},
    {'text': 'Falangas 1994 benchmark: -10.5 dB  (rate feedback on flat plate + PZT)', 'bold': True},
    {'text': ''},
    {'text': 'Gain sweep at 10%, 20%, 30%, 50%, 70%, 80% of g_max:'},
    {'text': 'Available in dvf_main.m (Figure 5 output)', 'level': 1},
    {'text': ''},
    {'text': 'Actuator voltage check: Va_peak = g * K_norm * H21(f1) * a_rms'},
    {'text': 'Limit: 400 V (P-876.A12 max).  Both stability and voltage limit checked.', 'level': 1},
],
x=MRG, y=Inches(1.25), w=Inches(7.0), h=Inches(5.6), size=18)

add_callout_box(sl,
    "Reference figure to compare: Falangas 1994, Figure 7 -- closed-loop transmissibility\n"
    "of flat aluminium plate with PZT patches and rate feedback. Target: indent at f1 ~10 dB deep.",
    Inches(7.5), Inches(1.5), Inches(5.5), Inches(1.5))

add_bullet_block(sl, [
    {'text': 'Not yet available (needs COMSOL re-export):'},
    {'text': 'T11 = H11 + H12 * K / (1 + H22*K) * H21  (tip transmissibility)', 'level': 1},
    {'text': '  -> Re-export H11 and H12 at IMU corner from Study 2 / Study 1', 'level': 2},
],
x=Inches(7.5), y=Inches(3.5), w=Inches(5.5), h=Inches(1.5), size=16)

# ==============================
# SLIDE 13: Conclusions
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl,
    "Conclusions: DVF with correct sign is the viable controller; H11/H12 re-export completes the MIMO validation",
    "Conclusions")
footer(sl, slide_num=sn, total=total_slides)

conclusions = [
    {'text': 'Phase I -- PPF with dual DuraAct patches:'},
    {'text': 'Unconditionally stable due to collocated piezo phase bound', 'level': 1},
    {'text': '-22.6 dB predicted suppression at 214 Hz (matches Baken 2025: -22.8 dB)', 'level': 1, 'bold': True},
    {'text': ''},
    {'text': 'Phase II -- DVF with MEMS accelerometer:'},
    {'text': 'PPF and IRC both inapplicable with acceleration sensor', 'level': 1},
    {'text': 'DVF with +g (not -g) restores genuine velocity damping with solid.accZ data', 'level': 1, 'bold': True},
    {'text': '-10 to -15 dB suppression predicted at 161 Hz (target mode)', 'level': 1, 'bold': True},
    {'text': ''},
    {'text': 'Critical model decisions:'},
    {'text': 'Domain Spring Foundation: only BC consistent between Study 1 and Study 2', 'level': 1},
    {'text': '161 Hz is the reliable DVF target; 87 Hz has feedthrough-distorted phase', 'level': 1},
    {'text': 'PIC255 d31 and epsR must be overridden from COMSOL defaults', 'level': 1},
    {'text': ''},
    {'text': 'Remaining work: re-export H11 and H12 at IMU corner to complete T11 calculation'},
]
add_bullet_block(sl, conclusions, x=MRG, y=Inches(1.25), w=W - 2*MRG, h=Inches(5.6), size=18)

# ==============================
# SLIDE 14: References
# ==============================
sn += 1
sl = prs.slides.add_slide(BLANK)
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

slide_header(sl, "References", "")
footer(sl, slide_num=sn, total=total_slides)

refs = [
    "[1] E.T. Falangas, J.A. Dworak, S. Koshigoe, \"Controlling plate vibrations using PZT actuators,\" IEEE Control Syst. Mag., 14(4):34-41, 1994.",
    "[2] G. Gatti, M.J. Brennan, P. Gardonio, \"Active damping of a beam using a collocated accelerometer and PZT patch,\" J. Sound Vib., 303:798-813, 2007.",
    "[3] S.S. Aphale, A.J. Fleming, S.O.R. Moheimani, \"Integral resonant control of collocated smart structures,\" Smart Mater. Struct., 16(2):439-446, 2007.",
    "[4] M.J. Balas, \"Direct velocity feedback control of large space structures,\" J. Guidance Control, 2(3):252-253, 1979.",
    "[5] C.J. Goh and T.K. Caughey, \"On the stability problem caused by finite actuator dynamics,\" Int. J. Control, 41(3):787-802, 1985.",
    "[6] M. Baken, V. Gupta, B. Jansen, S.H. HosseinNia, \"Harnessing PZT shear actuators for vibration control in sandwich beams,\" arXiv:2506.21713v1, 2025.",
]
y_ref = Inches(1.2)
for ref in refs:
    add_textbox(sl, ref, MRG, y_ref, W - 2*MRG, Inches(0.65), font_size=16, colour=DARK)
    y_ref += Inches(0.7)

# --------------------------------------------------------------------------
out = r"C:\Users\isaak\Desktop\Comsol\Internship\LatestFiles\ImuActuator\presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
